# ppt_app/streamlit_app.py

import streamlit as st
from pptx import Presentation
import io

# PPT 파일에서 텍스트를 추출하는 함수
def extract_text_from_pptx(uploaded_file):
    presentation = Presentation(io.BytesIO(uploaded_file.read()))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return "\n".join(text)

# Streamlit 앱 메인 함수
def main():
    st.title("PPT 응답 페이지")
    st.write("PPT 파일을 업로드하고 분석해보세요!")
    
    # 파일 업로드 기능
    uploaded_file = st.file_uploader("PPT 파일을 선택하세요", type=["pptx"])
    if uploaded_file is not None:
        st.write("파일이 업로드되었습니다!")
        # PPT 파일에서 텍스트를 추출하여 표시
        text = extract_text_from_pptx(uploaded_file)
        
        # 텍스트 영역 크기 조정 (height 값을 키워서 영역 크기 조정)
        st.text_area("PPT 내용", text, height=400)  # 높이를 400으로 설정 (기본값보다 더 큼)

# Streamlit 앱 실행
if __name__ == "__main__":
    main()
